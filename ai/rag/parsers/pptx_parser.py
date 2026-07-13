"""
PPTX Parser using python-pptx
=============================
"""

from pptx import Presentation
from pptx.shapes.base import BaseShape

from ai.rag.parsers.base_parser import BaseParser
from ai.rag.parsers.schemas import (
    FileType,
    LocationType,
    BlockType,
    DocumentBlock,
)
from ai.rag.parsers.exceptions import CorruptedDocumentError, EmptyDocumentError


class PPTXParser(BaseParser):
    """python-pptx 기반 PPTX 파서"""

    def get_file_type(self) -> FileType:
        return FileType.PPTX

    def get_page_count(self) -> int | None:
        try:
            prs = Presentation(str(self.file_path))
            return len(prs.slides)
        except Exception:
            return None

    def _extract_text_from_shape(self, shape: BaseShape) -> str | None:
        """도형에서 텍스트 추출"""
        if not shape.has_text_frame:
            return None
        return shape.text_frame.text.strip()

    def _get_shape_order_key(self, shape: BaseShape) -> tuple:
        """
        도형 정렬 순서 결정 (화면상 읽기 순서 기준)

        먼저 top 좌표로 행 분류 후, 같은 행은 left 좌표로 정렬
        """
        top = shape.top if hasattr(shape, "top") else 0
        left = shape.left if hasattr(shape, "left") else 0

        # top을 기준으로 행 구분 (50pt 단위로 같은 행으로 처리)
        row = int(top // 50)
        return (row, left)

    def _classify_shape(self, shape: BaseShape) -> BlockType:
        """도형 유형 분류"""
        shape_type = shape.shape_type

        # 제목 자리 (크기 비율로 판단)
        if hasattr(shape, "width") and hasattr(shape, "height"):
            width = shape.width
            # 너비가 너무 크면 제목일 가능성
            if width > 6000000:  # 약 6.6cm 이상
                text = self._extract_text_from_shape(shape)
                if text and len(text) < 100:
                    return BlockType.TITLE

        # 도형 유형별 분류
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape_type == MSO_SHAPE_TYPE.TABLE:
            return BlockType.TABLE
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            return BlockType.IMAGE

        return BlockType.SHAPE

    def _table_to_text(self, shape: BaseShape) -> str:
        """표 도형을 텍스트로 변환"""
        if not shape.has_table:
            return ""

        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("\t".join(cells))
        return "\n".join(rows)

    def parse(self) -> DocumentExtractionResult:
        """PPTX 문서 파싱"""
        file_size = self.file_path.stat().st_size
        warnings: list[str] = []
        blocks: list[DocumentBlock] = []
        global_order = 0

        try:
            prs = Presentation(str(self.file_path))
        except Exception as e:
            raise CorruptedDocumentError(f"PPTX 파일을 열 수 없습니다: {e}")

        document_id = self.generate_document_id(self.file_path)
        slide_count = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            # 도형을 읽기 순서대로 정렬
            shapes_with_text = []
            for shape in slide.shapes:
                text = self._extract_text_from_shape(shape)
                if text:
                    shapes_with_text.append((shape, text))

            # top, left 좌표로 정렬
            shapes_with_text.sort(key=lambda x: self._get_shape_order_key(x[0]))

            # 정렬된 순서로 블록 추가
            for shape, text in shapes_with_text:
                block_type = self._classify_shape(shape)

                # 표인 경우 별도 처리
                if block_type == BlockType.TABLE and shape.has_table:
                    content = self._table_to_text(shape)
                else:
                    content = text

                if content.strip():
                    block_obj = DocumentBlock(
                        block_id=self.generate_block_id(
                            document_id,
                            LocationType.SLIDE,
                            slide_num,
                            global_order,
                        ),
                        block_type=block_type,
                        content=content,
                        location_type=LocationType.SLIDE,
                        location_number=slide_num,
                        order=global_order,
                        metadata={
                            "shape_name": shape.name if hasattr(shape, "name") else None,
                            "shape_type": str(shape.shape_type) if hasattr(shape, "shape_type") else None,
                        },
                    )
                    blocks.append(block_obj)
                    global_order += 1

        if len(blocks) == 0:
            raise EmptyDocumentError(
                "PPTX에서 텍스트를 추출할 수 없습니다. 빈 슬라이드이거나 손상된 파일일 수 있습니다."
            )

        return self.create_result(
            file_size=file_size,
            page_count=slide_count,
            blocks=blocks,
            warnings=warnings,
        )
