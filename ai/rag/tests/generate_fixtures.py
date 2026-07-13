"""
Generate Test Fixtures
=====================
테스트용 샘플 문서를 생성합니다.

Usage:
    python -m ai.rag.tests.generate_fixtures
"""

import sys
from pathlib import Path

# 프로젝트 루트 설정
project_root = Path(__file__).parent.parent.parent.parent
sys.path.insert(0, str(project_root))

# pdfkit 사용 시 (설치되어 있지 않으면 예외 처리)
try:
    import fitz  # PyMuPDF
except ImportError:
    print("PyMuPDF가 설치되어 있지 않습니다.")
    fitz = None

try:
    from docx import Document
except ImportError:
    print("python-docx가 설치되어 있지 않습니다.")
    Document = None

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx가 설치되어 있지 않습니다.")
    Presentation = None


def create_sample_pdf(output_path: Path) -> None:
    """샘플 PDF 생성"""
    if fitz is None:
        print("PyMuPDF 설치 필요: pip install pymupdf")
        return

    doc = fitz.open()

    # 1페이지
    page = doc.new_page(width=595, height=842)
    page.insert_text((50, 100), "사업계획서", fontsize=24)
    page.insert_text((50, 150), "회사명: AI 스타트업", fontsize=14)
    page.insert_text((50, 180), "작성일: 2026년 7월", fontsize=14)

    # 2페이지
    page = doc.new_page(width=595, height=842)
    page.insert_text((50, 100), "1. Executive Summary", fontsize=18)
    page.insert_text((50, 150), "우리는 AI 기술을 활용하여 교육 혁신을",
                      fontsize=12)
    page.insert_text((50, 170), "이끌고자 합니다. 개인 맞춤형 학습 플랫폼을",
                      fontsize=12)
    page.insert_text((50, 190), "통해 학습 효율성을 30% 향상시키는 것이 목표입니다.",
                      fontsize=12)

    # 3페이지
    page = doc.new_page(width=595, height=842)
    page.insert_text((50, 100), "2. 시장 분석", fontsize=18)
    page.insert_text((50, 150), "국내 교육市场规模은 30조 원이며,",
                      fontsize=12)
    page.insert_text((50, 170), "연평균 5% 성장하고 있습니다.",
                      fontsize=12)

    doc.save(str(output_path))
    doc.close()
    print(f"생성됨: {output_path}")


def create_sample_docx(output_path: Path) -> None:
    """샘플 DOCX 생성"""
    if Document is None:
        print("python-docx 설치 필요: pip install python-docx")
        return

    doc = Document()

    # 제목
    doc.add_heading("프로젝트 제안서", 0)

    # 문단
    doc.add_paragraph("이 제안서는 AI 기반 교육 플랫폼 개발에 관한 것입니다.")
    doc.add_paragraph("목표: 학습 효율성 30% 향상")

    # 표
    table = doc.add_table(rows=3, cols=2)
    table.style = "Light Grid Accent 1"
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "항목"
    hdr_cells[1].text = "내용"
    row1_cells = table.rows[1].cells
    row1_cells[0].text = "예산"
    row1_cells[1].text = "1억 원"
    row2_cells = table.rows[2].cells
    row2_cells[0].text = "기간"
    row2_cells[1].text = "12개월"

    # 목록
    doc.add_heading("주요 기능", level=2)
    doc.add_paragraph("AI 기반 학습 분석", style="List Bullet")
    doc.add_paragraph("개인 맞춤형 콘텐츠 추천", style="List Bullet")
    doc.add_paragraph("실시간 피드백 시스템", style="List Bullet")

    doc.save(str(output_path))
    print(f"생성됨: {output_path}")


def pt_to_emu(pt: int) -> int:
    """포인트(PT)를 EMU로 변환"""
    return pt * 12700


def create_sample_pptx(output_path: Path) -> None:
    """샘플 PPTX 생성"""
    if Presentation is None:
        print("python-pptx 설치 필요: pip install python-pptx")
        return

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000

    # 슬라이드 1: 제목
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    title = slide1.shapes.add_textbox(500000, 2000000, 8000000, 1000000)
    title.text_frame.text = "AI 교육 플랫폼 제안"
    title.text_frame.paragraphs[0].font.size = pt_to_emu(44)

    subtitle = slide1.shapes.add_textbox(500000, 3500000, 8000000, 500000)
    subtitle.text_frame.text = "2026년 7월"
    subtitle.text_frame.paragraphs[0].font.size = pt_to_emu(24)

    # 슬라이드 2: 주요 기능
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title2 = slide2.shapes.add_textbox(500000, 500000, 8000000, 800000)
    title2.text_frame.text = "주요 기능"
    title2.text_frame.paragraphs[0].font.size = pt_to_emu(36)

    content2 = slide2.shapes.add_textbox(500000, 1500000, 8000000, 4000000)
    content2.text_frame.text = "1. AI 기반 학습 분석\n2. 맞춤형 콘텐츠 추천\n3. 실시간 피드백"
    content2.text_frame.paragraphs[0].font.size = pt_to_emu(24)

    # 슬라이드 3: 예상 효과
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    title3 = slide3.shapes.add_textbox(500000, 500000, 8000000, 800000)
    title3.text_frame.text = "예상 효과"
    title3.text_frame.paragraphs[0].font.size = pt_to_emu(36)

    content3 = slide3.shapes.add_textbox(500000, 1500000, 8000000, 4000000)
    content3.text_frame.text = "• 학습 효율성 30% 향상\n• 만족도 90% 이상\n• 이탈률 20% 감소"
    content3.text_frame.paragraphs[0].font.size = pt_to_emu(24)

    prs.save(str(output_path))
    print(f"생성됨: {output_path}")


def create_empty_pdf(output_path: Path) -> None:
    """빈 PDF 생성 (텍스트 없음)"""
    if fitz is None:
        print("PyMuPDF 설치 필요: pip install pymupdf")
        return

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    doc.save(str(output_path))
    doc.close()
    print(f"생성됨: {output_path}")


def create_corrupted_file(output_path: Path) -> None:
    """손상된 파일 생성 (유효하지 않은 PDF 헤더)"""
    with open(output_path, "wb") as f:
        f.write(b"%PDF-1.0\n%\xe2\xe3\xcf\xd3\n")
        f.write(b"This is not a valid PDF content")
    print(f"생성됨: {output_path}")


def main():
    """모든 테스트 픽스처 생성"""
    fixtures_dir = Path(__file__).parent / "fixtures"
    fixtures_dir.mkdir(exist_ok=True)

    print("테스트 픽스처 생성 중...\n")

    create_sample_pdf(fixtures_dir / "sample.pdf")
    create_sample_docx(fixtures_dir / "sample.docx")
    create_sample_pptx(fixtures_dir / "sample.pptx")
    create_empty_pdf(fixtures_dir / "empty.pdf")
    create_corrupted_file(fixtures_dir / "corrupted.pdf")

    # 지원하지 않는 형식 (이미 존재하는 텍스트 파일)
    txt_path = fixtures_dir / "sample.txt"
    if not txt_path.exists():
        txt_path.write_text("이것은 테스트용 텍스트 파일입니다.")
        print(f"생성됨: {txt_path}")

    print("\n모든 픽스처 생성 완료!")


if __name__ == "__main__":
    main()
